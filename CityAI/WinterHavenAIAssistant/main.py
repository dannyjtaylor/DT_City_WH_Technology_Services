
import os
from dotenv import load_dotenv
from openai import OpenAI
import chromadb
import pdfplumber
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

# === Environment Setup ===
dotenv_path = os.path.join(os.path.dirname(__file__), ".env")
load_dotenv(dotenv_path)

api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError("❌ OPENAI_API_KEY not found. Check your .env file.")

client = OpenAI(api_key=api_key)

# === Vector DB Setup ===
chroma_client = chromadb.Client()
collection = chroma_client.get_or_create_collection(name="city_docs")

SUPPORTED_EXTENSIONS = [".txt", ".docx", ".pptx", ".pdf"]
LOG_FILE = "embedding_log.txt"

# === File Readers ===
def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def load_and_split_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".txt":
        full_text = extract_text_from_txt(file_path)
    elif ext == ".docx":
        full_text = extract_text_from_docx(file_path)
    elif ext == ".pptx":
        full_text = extract_text_from_pptx(file_path)
    elif ext == ".pdf":
        full_text = extract_text_from_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_path}")

    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    docs = splitter.create_documents([full_text])
    return [doc.page_content for doc in docs]

def embed_documents(docs, filename):
    for i, chunk in enumerate(docs):
        embedding = client.embeddings.create(
            model="text-embedding-3-small",
            input=chunk
        ).data[0].embedding

        collection.add(
    documents=[f"{filename}\n\n{chunk}"],  # prepend filename as a tagk
    embeddings=[embedding],
    ids=[f"{filename}_chunk_{i}"]
)


def process_all_files_in_folder(folder_path):
    with open(LOG_FILE, "w", encoding="utf-8") as log:
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            ext = os.path.splitext(filename)[1].lower()

            if ext in SUPPORTED_EXTENSIONS:
                print(f"[INFO] Processing: {filename}")
                try:
                    chunks = load_and_split_text(file_path)
                    embed_documents(chunks, filename)
                    log.write(f"✅ Embedded: {filename}\n")
                except Exception as e:
                    log.write(f"❌ Failed: {filename} — {str(e)}\n")
                    print(f"[ERROR] Failed to process {filename}: {e}")
            else:
                log.write(f"⏩ Skipped: {filename} (unsupported)\n")
                print(f"[SKIPPED] Unsupported file: {filename}")

def ask_question(query):
    query_embedding = client.embeddings.create(
        model="text-embedding-3-small",
        input=query
    ).data[0].embedding

    results = collection.query(query_embeddings=[query_embedding], n_results=10)
    context = "\n".join(results["documents"][0])

    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a helpful assistant who only uses the provided context."},
            {"role": "user", "content": f"Use this info to answer:\n\n{context}\n\nQuestion: {query}"}
        ]
    )
    return response.choices[0].message.content

if __name__ == "__main__":
    BASE_DIR = os.path.dirname(__file__)
    data_folder = os.path.join(BASE_DIR, "data")

    print(f"Scanning files in: {data_folder}")
    process_all_files_in_folder(data_folder)

    print("\nHello! I'm Vira, your personal assistant for navigating your way through your job at the City of Winter Haven. Ask me anything! ('exit' to quit)\n")
    while True:
        user_input = input("You: ")
        if user_input.lower() in ["exit", "quit"]:
            break
        answer = ask_question(user_input)
        print(f"\nVira: {answer}\n")
