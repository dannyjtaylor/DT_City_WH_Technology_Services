
import os
from pathlib import Path
import json
import logging
from dotenv import load_dotenv
from fastapi import FastAPI, Request
from fastapi.responses import JSONResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from openai import OpenAI
import chromadb
from langchain.text_splitter import RecursiveCharacterTextSplitter
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

chroma_client = chromadb.Client()
collection = chroma_client.get_or_create_collection(name="city_docs")

def extract_text_from_file(file_path):
    ext = file_path.suffix.lower()
    if ext == ".txt":
        return Path(file_path).read_text(encoding="utf-8")
    elif ext == ".docx":
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    elif ext == ".pdf":
        reader = PdfReader(str(file_path))
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    elif ext == ".pptx":
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    return ""

def load_documents():
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    data_dir = Path("data")
    for file in data_dir.glob("*.*"):
        try:
            full_text = extract_text_from_file(file)
            docs = splitter.create_documents([full_text])
            for i, doc in enumerate(docs):
                embedding = client.embeddings.create(
                    model="text-embedding-3-small",
                    input=doc.page_content
                ).data[0].embedding
                collection.add(
                    documents=[doc.page_content],
                    embeddings=[embedding],
                    ids=[f"{file.name}_{i}"]
                )
        except Exception as e:
            print(f"[ERROR] Failed to process {file.name}: {e}")

load_documents()

app = FastAPI()
app.mount("/static", StaticFiles(directory="static"), name="static")

@app.get("/")
async def serve_index():
    return FileResponse("static/index.html")

@app.post("/ask")
async def ask_question(request: Request):
    body = await request.json()
    query = body.get("query", "")
    query_embedding = client.embeddings.create(
        model="text-embedding-3-small",
        input=query
    ).data[0].embedding

    results = collection.query(query_embeddings=[query_embedding], n_results=3)
    context = "\n".join(results['documents'][0]) if results['documents'] else ""
    prompt = f"Use this info:\n\n{context}\n\nQuestion: {query}"

    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a helpful assistant who only uses the provided context."},
            {"role": "user", "content": prompt}
        ]
    )
    answer = response.choices[0].message.content
    return JSONResponse(content={"answer": answer})
